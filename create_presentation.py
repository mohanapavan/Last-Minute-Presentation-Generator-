from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import requests
from io import BytesIO

def create_presentation(data):
    prs = Presentation()
    
    for i, slide_data in enumerate(data["slides"]):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        
        # Handle background image for first slide before anything else
        if i == 0 and "background_image" in slide_data:
            try:
                img = BytesIO(requests.get(slide_data["background_image"]).content)
                slide.shapes.add_picture(
                    img, 0, 0, 
                    width=prs.slide_width, 
                    height=prs.slide_height
                )
            except Exception as e:
                print(f"Background image error: {e}")

            # Manually add title textbox on top of background
            title_box = slide.shapes.add_textbox(
                left=Inches(1),
                top=Inches(0.5),
                width=prs.slide_width - Inches(2),
                height=Inches(1)
            )
            title_frame = title_box.text_frame
            title_frame.clear()
            p = title_frame.paragraphs[0]
            p.text = slide_data["title"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)  # White text
            p.alignment = PP_ALIGN.CENTER

        else:
            # Normal slides
            title = slide.shapes.title
            title.text = slide_data["title"]
            
            # Add content
            left = Inches(4)
            text_box = slide.shapes.add_textbox(
                left=left,
                top=Inches(1.5),
                width=Inches(6),
                height=Inches(4)
            )
            text_frame = text_box.text_frame
            for line in slide_data.get("content", []):
                p = text_frame.add_paragraph()
                p.text = line

            # Add side image
            if "side_image" in slide_data:
                try:
                    img = BytesIO(requests.get(slide_data["side_image"]).content)
                    slide.shapes.add_picture(
                        img, 
                        left=Inches(0.5), 
                        top=Inches(1.5),
                        width=prs.slide_width * 0.3,
                        height=prs.slide_height * 0.6
                    )
                except Exception as e:
                    print(f"Side image error: {e}")
    
    ppt_stream = BytesIO()
    prs.save(ppt_stream)
    ppt_stream.seek(0)
    return ppt_stream
